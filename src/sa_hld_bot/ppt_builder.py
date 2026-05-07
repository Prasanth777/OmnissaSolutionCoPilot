from __future__ import annotations

import re
from dataclasses import dataclass
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from .catalog import Product


FONT_PRIMARY = "Aptos"

BG = RGBColor(0xF5, 0xF8, 0xFC)
TITLE_BAND = RGBColor(0x0C, 0x2D, 0x57)
ACCENT = RGBColor(0x00, 0x7A, 0x7A)
TEXT = RGBColor(0x1D, 0x1D, 0x1D)
SUBTLE = RGBColor(0x5B, 0x66, 0x73)


def _clean_text(text: str, max_len: int = 420) -> str:
    raw = (text or "").strip()
    if not raw:
        return ""
    # Remove markdown markers and normalize whitespace.
    cleaned = re.sub(r"(^|\s)[#*_`>-]+", " ", raw)
    cleaned = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", cleaned)
    cleaned = re.sub(r"\s+", " ", cleaned).strip()
    return cleaned[:max_len]


@dataclass(frozen=True)
class TemplateProfile:
    title_left: int
    title_top: int
    title_width: int
    title_height: int
    title_font_name: str
    title_font_size_pt: float
    title_font_bold: bool
    title_font_rgb: tuple[int, int, int]
    image_left: int
    image_top: int
    image_width: int
    image_height: int
    slide_bg_rgb: tuple[int, int, int] | None = None


class HldPptBuilder:
    def build(
        self,
        output_path: Path,
        customer_name: str,
        selected_products: list[Product],
        questionnaire: dict[str, str],
        rag_narrative: dict[str, str],
        references: list[str],
        image_rows: list[dict[str, str]],
        sample_ppt_path: Path | None = None,
        use_sample_style: bool = False,
    ) -> Path:
        image_rows = [
            row
            for row in image_rows
            if row.get("image_type", "architecture_diagram") == "architecture_diagram"
            and Path(str(row.get("local_path", ""))).exists()
        ]
        profile = self._learn_from_sample(sample_ppt_path) if use_sample_style else None

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        self._title(prs, customer_name, selected_products, profile)

        for idx, row in enumerate(image_rows):
            self._image_only_slide(prs, row=row, idx=idx + 1, profile=profile)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)
        return output_path

    def _base_slide(self, prs: Presentation, title: str, profile: TemplateProfile | None = None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        fill = slide.background.fill
        fill.solid()
        if profile and profile.slide_bg_rgb:
            fill.fore_color.rgb = RGBColor(*profile.slide_bg_rgb)
        else:
            fill.fore_color.rgb = BG

        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.9))
        band.fill.solid()
        band.fill.fore_color.rgb = TITLE_BAND
        band.line.fill.background()

        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.9), Inches(13.333), Inches(0.05))
        accent.fill.solid()
        accent.fill.fore_color.rgb = ACCENT
        accent.line.fill.background()

        if profile:
            title_box = slide.shapes.add_textbox(
                profile.title_left,
                profile.title_top,
                profile.title_width,
                profile.title_height,
            )
        else:
            title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.14), Inches(12.6), Inches(0.58))
        p = title_box.text_frame.paragraphs[0]
        p.text = _clean_text(title, max_len=140)
        if profile:
            p.font.name = profile.title_font_name or FONT_PRIMARY
            p.font.size = Pt(max(14, profile.title_font_size_pt))
            p.font.bold = profile.title_font_bold
            p.font.color.rgb = RGBColor(*profile.title_font_rgb)
        else:
            p.font.name = FONT_PRIMARY
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        return slide

    def _title(
        self,
        prs: Presentation,
        customer_name: str,
        selected_products: list[Product],
        profile: TemplateProfile | None,
    ) -> None:
        slide = self._base_slide(prs, f"{customer_name or 'Customer'} | Architecture Design", profile)
        text = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(11.8), Inches(2.4)).text_frame
        p = text.paragraphs[0]
        p.text = "Architecture Diagram Set"
        p.font.name = FONT_PRIMARY
        p.font.size = Pt(19)
        p.font.color.rgb = TEXT
        p = text.add_paragraph()
        p.text = "Products: " + ", ".join(product.title for product in selected_products)
        p.font.name = FONT_PRIMARY
        p.font.size = Pt(15)
        p.font.color.rgb = SUBTLE

    def _image_only_slide(
        self,
        prs: Presentation,
        row: dict[str, str],
        idx: int,
        profile: TemplateProfile | None,
    ) -> None:
        title = row.get("slide_title") or row.get("title", f"Architecture Diagram {idx}")
        slide = self._base_slide(prs, title, profile)
        image_path = Path(str(row.get("local_path", "")))
        if image_path.exists():
            if profile:
                self._add_image_contain(
                    slide,
                    image_path,
                    profile.image_left,
                    profile.image_top,
                    profile.image_width,
                    profile.image_height,
                )
            else:
                self._add_image_contain(slide, image_path, Inches(0.55), Inches(1.1), Inches(12.25), Inches(5.85))

    @staticmethod
    def _context_bullets(questionnaire: dict[str, str]) -> list[str]:
        fields = [
            "customer_name",
            "industry",
            "project_scope",
            "users_personas",
            "hosting_strategy",
            "identity_source",
            "network_constraints",
            "security_requirements",
            "availability_requirements",
            "timeline",
            "assumptions",
        ]
        return [f"{field.replace('_', ' ').title()}: {questionnaire.get(field, 'TBD')}" for field in fields]

    def _add_image_contain(self, slide, image_path: Path, x, y, box_w, box_h) -> None:
        with Image.open(image_path) as img:
            w_px, h_px = img.size
        if w_px <= 0 or h_px <= 0:
            slide.shapes.add_picture(str(image_path), x, y, width=box_w)
            return
        image_ratio = w_px / h_px
        box_ratio = float(box_w) / float(box_h)
        if image_ratio >= box_ratio:
            width = box_w
            height = int(float(box_w) / image_ratio)
            top = int(float(y) + (float(box_h) - height) / 2.0)
            slide.shapes.add_picture(str(image_path), int(float(x)), top, width=int(float(width)))
        else:
            height = box_h
            width = int(float(box_h) * image_ratio)
            left = int(float(x) + (float(box_w) - width) / 2.0)
            slide.shapes.add_picture(str(image_path), left, int(float(y)), height=int(float(height)))

    def _learn_from_sample(self, sample_ppt_path: Path | None) -> TemplateProfile | None:
        if not sample_ppt_path or not sample_ppt_path.exists():
            return None
        try:
            sample = Presentation(sample_ppt_path)
        except Exception:
            return None

        best_slide = None
        best_picture = None
        for slide in sample.slides:
            pictures = [shape for shape in slide.shapes if shape.shape_type == 13]
            if not pictures:
                continue
            candidate = max(pictures, key=lambda shp: int(shp.width) * int(shp.height))
            if not best_picture or (int(candidate.width) * int(candidate.height) > int(best_picture.width) * int(best_picture.height)):
                best_picture = candidate
                best_slide = slide

        if best_slide is None or best_picture is None:
            return None

        title_shape = None
        for shape in best_slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            text = (shape.text or "").strip()
            if text and int(shape.top) < int(best_picture.top):
                title_shape = shape
                break
        if title_shape is None:
            return None

        para = title_shape.text_frame.paragraphs[0] if title_shape.text_frame.paragraphs else None
        run = para.runs[0] if para and para.runs else None
        font = run.font if run else (para.font if para else None)
        font_name = getattr(font, "name", None) or FONT_PRIMARY
        font_size_pt = float(getattr(font.size, "pt", 24) or 24)
        font_bold = bool(getattr(font, "bold", True))

        font_rgb = (255, 255, 255)
        try:
            rgb = getattr(getattr(font, "color", None), "rgb", None)
            if rgb is not None:
                font_rgb = (int(rgb[0]), int(rgb[1]), int(rgb[2]))
        except Exception:
            pass

        bg_rgb = None
        try:
            fill = best_slide.background.fill
            if fill and fill.type is not None and getattr(fill, "fore_color", None) and getattr(fill.fore_color, "rgb", None):
                rgb = fill.fore_color.rgb
                bg_rgb = (int(rgb[0]), int(rgb[1]), int(rgb[2]))
        except Exception:
            bg_rgb = None

        return TemplateProfile(
            title_left=int(title_shape.left),
            title_top=int(title_shape.top),
            title_width=int(title_shape.width),
            title_height=int(title_shape.height),
            title_font_name=font_name,
            title_font_size_pt=font_size_pt,
            title_font_bold=font_bold,
            title_font_rgb=font_rgb,
            image_left=int(best_picture.left),
            image_top=int(best_picture.top),
            image_width=int(best_picture.width),
            image_height=int(best_picture.height),
            slide_bg_rgb=bg_rgb,
        )
